#!/usr/bin/env python3
"""
augur.py — Augur AI Agent
Fetches RSS feeds, analyzes via Claude API, outputs a PowerPoint briefing.

Usage:
    python augur.py --feeds urls.md
    python augur.py --feeds urls.md --max-age 24 --output my_briefing.pptx
    python augur.py --feeds-url https://raw.githubusercontent.com/tximista64/Augur/main/urls.md

Requirements:
    pip install feedparser anthropic python-pptx
"""

import argparse
import json
import os
import re
import sys
import urllib.request
from datetime import datetime, timedelta, timezone
from pathlib import Path

try:
    import feedparser
except ImportError:
    print("[!] pip install feedparser", file=sys.stderr); sys.exit(1)

try:
    import anthropic
except ImportError:
    print("[!] pip install anthropic", file=sys.stderr); sys.exit(1)

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("[!] pip install python-pptx", file=sys.stderr); sys.exit(1)


# ─── CONFIG ────────────────────────────────────────────────────────────────────

MAX_ARTICLES   = 100
MAX_PER_FEED   = 3
MAX_AGE_HOURS  = 48
MAX_DESC_CHARS = 500
MODEL          = "claude-opus-4-7"

PROFILES = {
    "pentest":  ["Offensive", "Bug Bounty", "CTF", "Persons of Interest"],
    "blueteam": ["Blueteam", "Conf", "Persons of Interest"],
    "geopo":    ["Osint", "Geopo", "Podcasts"],
    "all":      None,
}

# ─── COLORS (dark theme) ───────────────────────────────────────────────────────

C_BG     = RGBColor(0x0D, 0x11, 0x17)
C_PANEL  = RGBColor(0x16, 0x1B, 0x22)
C_BORDER = RGBColor(0x30, 0x36, 0x3D)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_GRAY   = RGBColor(0x8B, 0x94, 0x9E)
C_BLUE   = RGBColor(0x58, 0xA6, 0xFF)
C_RED    = RGBColor(0xF8, 0x51, 0x49)
C_ORANGE = RGBColor(0xFF, 0x7B, 0x0E)
C_YELLOW = RGBColor(0xE3, 0xB3, 0x41)
C_GREEN  = RGBColor(0x3F, 0xB9, 0x50)

CRIT_COLORS = {"🔴": C_RED, "🟠": C_ORANGE, "🟡": C_YELLOW, "🟢": C_GREEN}

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ─── IOC REGEX ─────────────────────────────────────────────────────────────────

_IOC = {
    "CVE":    re.compile(r'\bCVE-\d{4}-\d{4,7}\b'),
    "SHA256": re.compile(r'\b[0-9a-fA-F]{64}\b'),
    "MD5":    re.compile(r'\b[0-9a-fA-F]{32}\b'),
    "IP":     re.compile(r'\b(?:(?:25[0-5]|2[0-4]\d|[01]?\d\d?)\.){3}(?:25[0-5]|2[0-4]\d|[01]?\d\d?)\b'),
    "Domain": re.compile(r'\b(?:[a-zA-Z0-9-]{1,63}\.)+(?:ru|cn|ir|kp|onion|xyz|top|tk|ml)\b', re.I),
}
_PRIVATE = re.compile(r'^(?:127\.|10\.|172\.(?:1[6-9]|2\d|3[01])\.|192\.168\.)')
IOC_ORDER = {"CVE": 0, "IP": 1, "Domain": 2, "SHA256": 3, "MD5": 4}


def extract_iocs_regex(articles: list[dict]) -> list[dict]:
    seen, result = set(), []
    for a in articles:
        text = f"{a['title']} {a['desc']}"
        for ioc_type, pat in _IOC.items():
            for val in pat.findall(text):
                if ioc_type == "IP" and _PRIVATE.match(val):
                    continue
                key = (ioc_type, val.lower())
                if key in seen:
                    continue
                seen.add(key)
                crit = "🔴" if ioc_type in ("SHA256", "MD5", "Domain") else "🟠"
                result.append({
                    "type": ioc_type,
                    "value": val,
                    "severity": crit,
                    "source": a["source"],
                    "context": a["title"][:55],
                })
    return result[:30]


# ─── FEED LOADING ──────────────────────────────────────────────────────────────

def load_feeds(path=None, url=None, profile="all") -> list[str]:
    filt = PROFILES.get(profile)
    if url:
        with urllib.request.urlopen(url) as r:
            content = r.read().decode("utf-8")
    else:
        content = Path(path).read_text(encoding="utf-8")
    urls, section = [], ""
    for line in content.splitlines():
        s = line.strip()
        if not s:
            continue
        if s.startswith("#"):
            section = s.lstrip("#").strip()
            continue
        if filt and not any(f.lower() in section.lower() for f in filt):
            continue
        m = re.search(r'\(https?://[^\)]+\)', s)
        if m:
            urls.append(m.group()[1:-1])
        elif s.startswith("http"):
            urls.append(s)
    return list(dict.fromkeys(urls))


def is_recent(entry, max_age_hours: int) -> bool:
    for attr in ("published_parsed", "updated_parsed"):
        t = getattr(entry, attr, None)
        if t:
            pub = datetime(*t[:6], tzinfo=timezone.utc)
            return datetime.now(timezone.utc) - pub < timedelta(hours=max_age_hours)
    return True


def fetch_articles(feed_urls: list[str], max_age_hours: int) -> list[dict]:
    articles = []
    for url in feed_urls:
        try:
            feed = feedparser.parse(url, request_headers={"User-Agent": "augur/2.0"})
            source = feed.feed.get("title", url)
            count = 0
            for entry in feed.entries:
                if count >= MAX_PER_FEED:
                    break
                if not is_recent(entry, max_age_hours):
                    continue
                title = entry.get("title", "—")
                desc = re.sub(r"<[^>]+>", "", entry.get("summary", entry.get("description", "")))
                articles.append({
                    "source": source,
                    "title": title,
                    "desc": desc[:MAX_DESC_CHARS],
                    "link": entry.get("link", ""),
                })
                count += 1
        except Exception as e:
            print(f"  [!] {url}: {e}", file=sys.stderr)
    return articles[:MAX_ARTICLES]


# ─── CLAUDE ANALYSIS ───────────────────────────────────────────────────────────

_PROMPT = """You are a senior CTI analyst. Analyze the following cybersecurity articles and return a structured JSON object for a PowerPoint threat intelligence briefing.

Date: {date}

RETURN ONLY VALID JSON — no markdown fences, no commentary.

{{
  "exec_summary": [
    "<one-sentence top threat #1 — most impactful finding>",
    "<one-sentence top threat #2>",
    "<one-sentence top threat #3>"
  ],
  "cve_vulns": [
    {{
      "title": "<concise title>",
      "severity": "🔴|🟠|🟡|🟢",
      "summary": "<2 sentences, SOC/DFIR-oriented, factual>",
      "source": "<feed name>",
      "link": "<url>"
    }}
  ],
  "apt_actors":   [{{...same fields...}}],
  "offensive_it": [{{...same fields...}}],
  "ransomware":   [{{...same fields...}}],
  "geopolitics":  [{{...same fields...}}],
  "iocs_claude": [
    {{
      "type": "IP|Domain|Hash|CVE",
      "value": "<exact indicator>",
      "severity": "🔴|🟠|🟡|🟢",
      "source": "<feed name>",
      "context": "<10 words max describing threat context>"
    }}
  ]
}}

Classification rules:
- cve_vulns    : CVEs, patches, vulnerabilities, exploits, 0days
- apt_actors   : nation-state groups, campaigns, TTPs, attribution
- offensive_it : red team tools, techniques, pentest, post-exploitation, C2, shellcode
- ransomware   : ransomware groups, extortion, leaks, cybercrime operations
- geopolitics  : geopolitical cyber context, OSINT, strategic intel, sanctions, conflicts

Scoring rules:
- 🔴 Critical : actively exploited, RCE, 0day, confirmed APT intrusion
- 🟠 High     : public PoC, high-impact vuln, active campaign
- 🟡 Medium   : patched vuln, research, no active exploitation
- 🟢 Info     : general awareness, strategic context

Output rules:
- Each section: max 5 items, sorted by criticality descending
- exec_summary: 3 most impactful cross-section findings
- iocs_claude: only IOCs explicitly mentioned in article text (not inferred)
- Empty section → return []

ARTICLES:
{articles}"""


def analyze(articles: list[dict], date_str: str) -> dict:
    articles_txt = ""
    for i, a in enumerate(articles, 1):
        articles_txt += (
            f"\n[{i}] {a['source']}\n"
            f"Title: {a['title']}\n"
            f"Summary: {a['desc']}\n"
            f"Link: {a['link']}\n"
        )

    client = anthropic.Anthropic()
    print(f"⟳ Calling Claude API ({MODEL})...", file=sys.stderr)

    response = client.messages.create(
        model=MODEL,
        max_tokens=4096,
        messages=[{
            "role": "user",
            "content": _PROMPT.format(date=date_str, articles=articles_txt),
        }],
    )

    raw = response.content[0].text.strip()
    m = re.search(r'```(?:json)?\s*([\s\S]+?)\s*```', raw)
    if m:
        raw = m.group(1)

    return json.loads(raw)


# ─── PPTX PRIMITIVES ───────────────────────────────────────────────────────────

def _blank_layout(prs):
    for layout in prs.slide_layouts:
        if "Blank" in layout.name:
            return layout
    return prs.slide_layouts[-1]


def _set_bg(slide, color=None):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color or C_BG


def _rect(slide, left, top, width, height, fill, line=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.5)
    else:
        shape.line.fill.background()
    return shape


def _text(slide, text, left, top, width, height,
          size=14, bold=False, color=None, align=PP_ALIGN.LEFT):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color or C_WHITE
    run.font.name = "Calibri"
    return txb


def _accent_bars(slide, color):
    _rect(slide, 0, 0, SLIDE_W, Inches(0.07), color)
    _rect(slide, 0, 0, Inches(0.06), SLIDE_H, color)


# ─── SLIDE BUILDERS ────────────────────────────────────────────────────────────

def _slide_title(prs, date_str, n_articles, n_sources):
    slide = prs.slides.add_slide(_blank_layout(prs))
    _set_bg(slide)
    _rect(slide, 0, 0, SLIDE_W, Inches(0.07), C_RED)
    _rect(slide, 0, Inches(7.43), SLIDE_W, Inches(0.07), C_BLUE)

    _text(slide, "AUGUR",
          Inches(1), Inches(1.8), Inches(11.333), Inches(1.6),
          size=58, bold=True, align=PP_ALIGN.CENTER)
    _text(slide, "THREAT INTELLIGENCE BRIEFING",
          Inches(1), Inches(3.3), Inches(11.333), Inches(0.7),
          size=18, color=C_BLUE, align=PP_ALIGN.CENTER)
    _text(slide, date_str,
          Inches(1), Inches(4.05), Inches(11.333), Inches(0.6),
          size=14, color=C_GRAY, align=PP_ALIGN.CENTER)
    _text(slide, f"{n_articles} articles  ·  {n_sources} sources",
          Inches(1), Inches(4.75), Inches(11.333), Inches(0.5),
          size=11, color=C_GRAY, align=PP_ALIGN.CENTER)


def _slide_exec_summary(prs, bullets: list[str]):
    slide = prs.slides.add_slide(_blank_layout(prs))
    _set_bg(slide)
    _accent_bars(slide, C_BLUE)

    _text(slide, "EXECUTIVE SUMMARY",
          Inches(0.3), Inches(0.12), Inches(12), Inches(0.5),
          size=10, bold=True, color=C_BLUE)
    _text(slide, "Top Threats Today",
          Inches(0.3), Inches(0.6), Inches(12), Inches(0.9),
          size=26, bold=True)

    for i, bullet in enumerate(bullets[:3]):
        y = Inches(1.65 + i * 1.65)
        _rect(slide, Inches(0.3), y, Inches(12.7), Inches(1.45), C_PANEL)
        _text(slide, ["①", "②", "③"][i],
              Inches(0.45), y + Inches(0.1), Inches(0.6), Inches(1.2),
              size=26, color=C_BLUE, bold=True)
        _text(slide, bullet,
              Inches(1.15), y + Inches(0.2), Inches(11.5), Inches(1.0),
              size=15, color=C_WHITE)


def _slide_section(prs, title: str, icon: str, items: list[dict], accent: RGBColor):
    if not items:
        return
    slide = prs.slides.add_slide(_blank_layout(prs))
    _set_bg(slide)
    _accent_bars(slide, accent)

    _text(slide, f"{icon}  {title.upper()}",
          Inches(0.3), Inches(0.12), Inches(12), Inches(0.5),
          size=10, bold=True, color=accent)

    item_h = Inches(1.22)
    gap    = Inches(0.12)
    y0     = Inches(0.68)

    for i, item in enumerate(items[:5]):
        y = y0 + i * (item_h + gap)
        ccolor = CRIT_COLORS.get(item.get("severity", "🟢"), C_GREEN)

        _rect(slide, Inches(0.3),  y, Inches(12.7), item_h, C_PANEL)
        _rect(slide, Inches(0.3),  y, Inches(0.06), item_h, ccolor)

        _text(slide, item.get("severity", ""),
              Inches(0.42), y + Inches(0.06), Inches(0.4), Inches(0.4), size=13)
        _text(slide, item.get("title", ""),
              Inches(0.92), y + Inches(0.06), Inches(11.7), Inches(0.45),
              size=13, bold=True)
        _text(slide, item.get("summary", ""),
              Inches(0.92), y + Inches(0.52), Inches(10.5), Inches(0.6),
              size=10, color=C_GRAY)
        _text(slide, item.get("source", ""),
              Inches(11.2), y + Inches(0.52), Inches(1.8), Inches(0.6),
              size=9, color=C_BLUE, align=PP_ALIGN.RIGHT)


def _slide_iocs(prs, iocs: list[dict]):
    if not iocs:
        return
    slide = prs.slides.add_slide(_blank_layout(prs))
    _set_bg(slide)
    _accent_bars(slide, C_RED)

    _text(slide, "🔍  IOCs OF THE DAY",
          Inches(0.3), Inches(0.12), Inches(12), Inches(0.5),
          size=10, bold=True, color=C_RED)

    iocs_sorted = sorted(iocs, key=lambda x: (IOC_ORDER.get(x["type"], 9), x.get("severity", "")))

    _text(slide, f"{len(iocs_sorted)} indicators  ·  regex + Claude extraction",
          Inches(0.3), Inches(0.6), Inches(12), Inches(0.35),
          size=9, color=C_GRAY)

    ROW_H = Inches(0.40)
    y_hdr = Inches(1.02)
    cols   = [Inches(0.30), Inches(1.55), Inches(6.35), Inches(9.55), Inches(11.75)]
    widths = [Inches(1.20), Inches(4.75), Inches(3.15), Inches(2.15), Inches(1.55)]
    hdrs   = ["TYPE", "INDICATOR", "CONTEXT", "SOURCE", "CRIT."]

    _rect(slide, Inches(0.3), y_hdr, Inches(12.7), ROW_H, C_BORDER)
    for hdr, col, w in zip(hdrs, cols, widths):
        _text(slide, hdr, col, y_hdr + Inches(0.07), w, Inches(0.28),
              size=8, bold=True, color=C_GRAY)

    for i, ioc in enumerate(iocs_sorted[:12]):
        y = Inches(1.42) + i * (ROW_H + Inches(0.03))
        _rect(slide, Inches(0.3), y, Inches(12.7), ROW_H, C_PANEL if i % 2 == 0 else C_BG)

        ccolor = CRIT_COLORS.get(ioc.get("severity", "🟢"), C_GREEN)
        _rect(slide, Inches(0.3), y, Inches(0.04), ROW_H, ccolor)

        val = ioc["value"]
        if ioc["type"] in ("SHA256", "MD5") and len(val) > 22:
            val = val[:10] + "…" + val[-10:]

        for v, col, w, color in zip(
            [ioc["type"], val, ioc.get("context", "")[:55], ioc.get("source", "")[:28], ioc.get("severity", "")],
            cols, widths,
            [C_GRAY, C_WHITE, C_GRAY, C_GRAY, C_GRAY],
        ):
            _text(slide, v, col, y + Inches(0.07), w, Inches(0.28), size=8, color=color)


def _slide_stats(prs, n_articles, n_sources, n_iocs, date_str):
    slide = prs.slides.add_slide(_blank_layout(prs))
    _set_bg(slide)
    _rect(slide, 0, 0, SLIDE_W, Inches(0.07), C_BLUE)

    _text(slide, "STATS & SOURCES",
          Inches(0.3), Inches(0.12), Inches(12), Inches(0.5),
          size=10, bold=True, color=C_BLUE)
    _text(slide, date_str,
          Inches(0.3), Inches(0.7), Inches(12), Inches(0.5),
          size=13, color=C_GRAY)

    for i, (num, label) in enumerate([
        (str(n_articles), "Articles analyzed"),
        (str(n_sources),  "RSS sources"),
        (str(n_iocs),     "IOCs extracted"),
    ]):
        x = Inches(1.3 + i * 3.9)
        _rect(slide, x, Inches(1.9), Inches(3.4), Inches(2.8), C_PANEL)
        _text(slide, num, x, Inches(2.1), Inches(3.4), Inches(1.4),
              size=52, bold=True, color=C_BLUE, align=PP_ALIGN.CENTER)
        _text(slide, label, x, Inches(3.45), Inches(3.4), Inches(0.7),
              size=13, color=C_GRAY, align=PP_ALIGN.CENTER)

    _text(slide, "Augur AI Agent  ·  github.com/tximista64/Augur",
          Inches(0.3), Inches(7.1), Inches(12.7), Inches(0.35),
          size=8, color=C_GRAY, align=PP_ALIGN.CENTER)


# ─── BUILD ─────────────────────────────────────────────────────────────────────

def build_pptx(data: dict, iocs_regex: list[dict], n_articles: int,
               n_sources: int, date_str: str, output: str):
    seen, iocs_all = set(), []
    for ioc in iocs_regex + data.get("iocs_claude", []):
        key = (ioc["type"], ioc["value"].lower())
        if key not in seen:
            seen.add(key)
            iocs_all.append(ioc)

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    _slide_title(prs, date_str, n_articles, n_sources)
    _slide_exec_summary(prs, data.get("exec_summary", []))
    _slide_section(prs, "CVE & Vulnerabilities",   "🛡️", data.get("cve_vulns",    []), C_RED)
    _slide_section(prs, "APT & Threat Actors",     "🎯", data.get("apt_actors",   []), C_ORANGE)
    _slide_section(prs, "Offensive IT / Red Team", "⚔️", data.get("offensive_it", []), C_YELLOW)
    _slide_section(prs, "Ransomware & Cybercrime", "💀", data.get("ransomware",   []), C_ORANGE)
    _slide_section(prs, "Geopolitics & Intel",     "🌍", data.get("geopolitics",  []), C_BLUE)
    _slide_iocs(prs, iocs_all)
    _slide_stats(prs, n_articles, n_sources, len(iocs_all), date_str)

    prs.save(output)
    print(f"✓ {output}  ({len(prs.slides)} slides · {len(iocs_all)} IOCs)", file=sys.stderr)


# ─── MAIN ──────────────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(
        description="Augur AI Agent — PowerPoint briefing from RSS feeds"
    )
    g = parser.add_mutually_exclusive_group(required=True)
    g.add_argument("--feeds",     metavar="FILE", help="Local .md feed list")
    g.add_argument("--feeds-url", metavar="URL",  help="Remote feed list URL")
    parser.add_argument("--max-age", type=int, default=MAX_AGE_HOURS, metavar="H",
                        help=f"Max article age in hours (default: {MAX_AGE_HOURS})")
    parser.add_argument("--profile", choices=list(PROFILES), default="all",
                        help="Feed profile filter (default: all)")
    parser.add_argument("--output", metavar="FILE", default=None,
                        help="Output .pptx path (default: briefing_YYYYMMDD.pptx)")
    args = parser.parse_args()

    if not os.getenv("ANTHROPIC_API_KEY"):
        print("[!] ANTHROPIC_API_KEY not found in environment.", file=sys.stderr)
        print("    Add it to your ~/.zshrc:  export ANTHROPIC_API_KEY=sk-ant-...", file=sys.stderr)
        sys.exit(1)

    date_str = datetime.now().strftime("%d/%m/%Y")
    output   = args.output or f"briefing_{datetime.now().strftime('%Y%m%d')}.pptx"

    print(f"⟳ Loading feeds [profile: {args.profile}]...", file=sys.stderr)
    feed_urls = load_feeds(path=args.feeds, url=args.feeds_url, profile=args.profile)
    if not feed_urls:
        print("[!] No feed URLs found.", file=sys.stderr); sys.exit(1)
    print(f"  → {len(feed_urls)} feeds", file=sys.stderr)

    print(f"⟳ Fetching articles (< {args.max_age}h)...", file=sys.stderr)
    articles = fetch_articles(feed_urls, args.max_age)
    if not articles:
        print("[!] No recent articles found.", file=sys.stderr); sys.exit(1)
    print(f"  → {len(articles)} articles", file=sys.stderr)

    iocs_regex = extract_iocs_regex(articles)
    print(f"  → {len(iocs_regex)} IOCs (regex)", file=sys.stderr)

    data = analyze(articles, date_str)
    print(f"  → {len(data.get('iocs_claude', []))} IOCs (Claude)", file=sys.stderr)

    build_pptx(data, iocs_regex, len(articles), len(feed_urls), date_str, output)


if __name__ == "__main__":
    main()
